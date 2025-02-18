import os
import json
import argparse
import re  # already imported globally
from urllib.parse import urlparse
import threading # New import

# For Azure Speech SDK
import azure.cognitiveservices.speech as speechsdk
import requests
from bs4 import BeautifulSoup

# New global import for Google TTS.
from google.cloud import texttospeech
from concurrent.futures import ThreadPoolExecutor
import unicodedata

def load_config():
    config_path = os.path.expanduser("~/.local/tts-cli/config.json")
    with open(config_path, "r") as f:
        return json.load(f)

def clean_markdown(text):
    text = re.sub(r'\n', ' ', text)
    text = re.sub(r'([*_~`#>\-]|^[\d+\. ])', '', text)
    text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
    text = re.sub(r'!\[(.*?)\]\(.*?\)', r'\1', text)
    text = re.sub(r'<(.*?)>', '', text)
    text = re.sub(r'\{.*\}', '', text)
    text = re.sub(r' +', ' ', text)
    return text

# New: clean text to keep only regular speakable characters
def clean_speakable_text(text):
    # Normalize text to decompose combined characters.
    text = unicodedata.normalize('NFKD', text)
    # Remove non-ASCII characters and keep letters, numbers, spaces and common punctuation.
    text = re.sub(r"[^a-zA-Z0-9\s\.\,\?\!\:\;\-']", '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def read_text_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext in ['.txt', '.md']:
        with open(filepath, "r", encoding="utf-8") as f:
            text = f.read()
        if ext == '.md':
            text = clean_markdown(text)
        return text
    elif ext == '.pdf':
        from PyPDF2 import PdfReader
        text = ""
        with open(filepath, "rb") as f:
            reader = PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
        # Remove wrapped lines to form continuous sentences
        text = text.replace("-\n", "").replace("\n", " ")
        return text
    elif ext == '.doc':
        import subprocess
        result = subprocess.run(["antiword", filepath], stdout=subprocess.PIPE)
        return result.stdout.decode("utf-8")
    elif ext == '.docx':
        import docx
        text = []
        doc = docx.Document(filepath)
        for para in doc.paragraphs:
            text.append(para.text)
        return "\n".join(text)
    elif ext in ['.ppt', '.pptx']:
        try:
            from pptx import Presentation
        except ImportError:
            raise Exception("Please install 'python-pptx' to support ppt/pptx files")
        prs = Presentation(filepath)
        slides_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides_text.append(shape.text)
        return "\n".join(slides_text)
    else:
        raise Exception("Unsupported file format!")

def extract_url_text(url):
    response = requests.get(url)
    response.raise_for_status()  # raise error for bad status
    soup = BeautifulSoup(response.content, "html.parser")
    article = soup.find("article")
    if article:
        return article.get_text(separator=" ", strip=True)
    main_tag = soup.find("main")
    if main_tag:
        return main_tag.get_text(separator=" ", strip=True)
    # Fallback: remove common non-article tags
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()
    return soup.body.get_text(separator=" ", strip=True)

def sanitize_filename(name):
    parsed = urlparse(name)
    combined = (parsed.netloc + parsed.path).strip("/")
    return re.sub(r"[^A-Za-z0-9_-]+", "_", combined)

def synthesize_azure(text, provider_config, output_file, force=False):
    speech_key = provider_config.get("speech_key")
    service_region = provider_config.get("service_region")
    voice_name = provider_config.get("voice")
    speech_config = speechsdk.SpeechConfig(subscription=speech_key, region=service_region)
    speech_config.speech_synthesis_voice_name = voice_name

    # Set output format if mp3 is required
    if provider_config.get("output_format", "wav") == "mp3":
        speech_config.set_speech_synthesis_output_format(
            speechsdk.SpeechSynthesisOutputFormat.Audio16Khz128KBitRateMonoMp3
        )

    # Split text if exceeds max_length instead of truncating
    max_length = provider_config.get("max_text_length", 5000)
    if len(text) > max_length:
        print(f"Text length ({len(text)}) exceeds {max_length}, splitting into chunks.")
        import textwrap
        # Split text into chunks without breaking words
        def chunk_text(text, max_length):
            return textwrap.wrap(text, width=max_length, break_long_words=False, break_on_hyphens=False)
        chunks = chunk_text(text, max_length)
        total = len(chunks)
        # Prompt confirmation if more than one chunk and force flag not set
        if total > 1 and not force:
            answer = input(f"This text will be split into {total} chunks. Proceed? (Y/n): ")
            if answer.strip().lower().startswith("n"):
                print("Aborting multi-chunk synthesis.")
                return
        base_audio = os.path.splitext(output_file)[0]
        chunk_files = []  # Track individual chunk audio files
        for i, chunk in enumerate(chunks, start=1):
            part_file = output_file.replace(".", f"_part{i}.")
            chunk_files.append(part_file)
            # Output each chunk's text to a file
            chunk_text_file = f"{base_audio}_part{i}-text.txt"
            with open(chunk_text_file, "w", encoding="utf-8") as f:
                f.write(chunk)
            print(f"Chunk {i} text saved to {chunk_text_file}")

            print(f"Synthesizing chunk {i}/{total} to {part_file}")
            audio_config = speechsdk.audio.AudioOutputConfig(filename=part_file)
            chunk_synth = speechsdk.SpeechSynthesizer(speech_config=speech_config, audio_config=audio_config)
            result = chunk_synth.speak_text_async(chunk).get()
            if result.reason != speechsdk.ResultReason.SynthesizingAudioCompleted:
                print(f"Error in chunk {i}: {getattr(result, 'error_details', 'No details available')}")
                raise Exception(f"Speech synthesis failed for chunk {i}")
            print(f"Chunk {i} saved to {part_file}")

        # Merge audio chunks using ffmpeg directly
        if len(chunk_files) > 1:
            import subprocess
            base_audio = os.path.splitext(output_file)[0]
            temp_list = f"{base_audio}_filelist.txt"
            with open(temp_list, "w", encoding="utf-8") as f:
                for file in chunk_files:
                    f.write(f"file '{os.path.abspath(file)}'\n")
            cmd = [
                "ffmpeg",
                "-y",
                "-f", "concat",
                "-safe", "0",
                "-i", temp_list,
                "-c", "copy",
                output_file
            ]
            subprocess.run(cmd, check=True)
            print(f"Merged audio saved to {output_file}")
            os.remove(temp_list)
            # Remove temporary chunk audio files and corresponding text files
            for i, audio_file in enumerate(chunk_files, start=1):
                try:
                    os.remove(audio_file)
                    # Removed logging for successful removal of audio file.
                except OSError as e:
                    print(f"Error removing file {audio_file}: {e}")
                text_file = f"{base_audio}_part{i}-text.txt"
                try:
                    os.remove(text_file)
                    # Removed logging for successful removal of text file.
                except OSError as e:
                    print(f"Error removing text file {text_file}: {e}")
    else:
        # Single chunk synthesis
        audio_config = speechsdk.audio.AudioOutputConfig(filename=output_file)
        synthesizer = speechsdk.SpeechSynthesizer(speech_config=speech_config, audio_config=audio_config)
        result = synthesizer.speak_text_async(text).get()
        if result.reason != speechsdk.ResultReason.SynthesizingAudioCompleted:
            print(f"Error details: {getattr(result, 'error_details', 'No details available')}")
            raise Exception("Speech synthesis failed")
        print(f"Audio saved to {output_file}")

# New: recursively process a sentence so that its byte length is under max_sentence_bytes.
def process_sentence(sentence, max_sentence_bytes):
    sentence = sentence.strip()
    if len(sentence.encode("utf-8")) <= max_sentence_bytes:
        return [sentence]
    words = sentence.split()
    pieces = []
    current = ""
    for word in words:
        candidate = (current + " " + word).strip() if current else word
        if len(candidate.encode("utf-8")) <= max_sentence_bytes:
            current = candidate
        else:
            pieces.append(current)
            current = word
    if current:
        pieces.append(current)
    return pieces

# In the synthesize_google function, update the sentence splitting logic:
def synthesize_google(text, provider_config, output_file, num_threads=1):
    max_bytes = 4000  # Google Cloud TTS limit in bytes
    max_sentence_bytes = 4000  # Maximum byte size for a single sentence
    base_audio = os.path.splitext(output_file)[0]
    output_ext = os.path.splitext(output_file)[1]
    state_file = f"{base_audio}_state.json"  # Define state file name
    state_lock = threading.Lock() # Create a lock

    def save_state(state):
        with open(state_file, "w", encoding="utf-8") as f:
            json.dump(state, f, indent=4)

    def load_state():
        try:
            with open(state_file, "r", encoding="utf-8") as f:
                return json.load(f)
        except FileNotFoundError:
            return None

    # Load state if it exists, otherwise initialize
    state = load_state()
    chunks = []
    if state is None:
        sentences = re.split(r'(?<=[.!?])\s+', text)

        final_sentences = []
        for sentence in sentences:
            pieces = process_sentence(sentence, max_sentence_bytes)
            # Join the pieces so that TTS treats them as multiple sentences.
            final_sentences.append(". ".join(pieces))

        current_chunk = ""
        current_chunk_bytes = 0

        for sentence in final_sentences:
            s_bytes = len(sentence.encode('utf-8'))
            if current_chunk_bytes + s_bytes <= max_bytes:
                current_chunk += (" " + sentence) if current_chunk else sentence
                current_chunk_bytes += s_bytes
            else:
                if current_chunk:
                    chunks.append(current_chunk)
                current_chunk = sentence
                current_chunk_bytes = s_bytes

        if current_chunk:
            chunks.append(current_chunk)

        state = {"completed_chunks": [], "chunks": chunks}
        save_state(state)  # Save the initial state to create the file
    else:
        chunks = state["chunks"]

    completed_chunks = state["completed_chunks"]
    total = len(chunks)

    # Create a list of chunks that haven't been processed yet
    remaining_chunks = [
        (i, chunk) for i, chunk in enumerate(chunks)
        if f"{base_audio}_part{i+1}{output_ext}" not in completed_chunks
    ]

    # Print information before starting synthesis
    total_chars = len(text)
    print(f"Total characters: {total_chars}")
    print(f"Number of batches: {total}")
    print(f"Remaining batches: {len(remaining_chunks)}")
    print(f"Completed: {len(completed_chunks)}/{total}")
    print("Progress: {:.2f}%".format(len(completed_chunks) / total * 100))

    def synthesize_chunk(i, chunk):
        part_file = f"{base_audio}_part{i+1}{output_ext}"

        if part_file in completed_chunks and os.path.exists(part_file):
            print(f"Chunk {i+1} already synthesized, skipping.")
            progress = (i + 1) / total * 100
            print(f"Progress: {progress:.2f}%")
            return  # Skip synthesis if file exists

        chunk_text_file = f"{base_audio}_part{i+1}-text.txt"
        with open(chunk_text_file, "w", encoding="utf-8") as f:
            f.write(chunk)
        print(f"Chunk {i+1} text saved to {chunk_text_file}")

        print(f"Synthesizing chunk {i+1}/{total} to {part_file}")

        try:
            client = texttospeech.TextToSpeechClient()
            voice_name = provider_config.get("voice", "en-US-Wavenet-D")
            language_code = "-".join(voice_name.split("-")[:2])
            voice = texttospeech.VoiceSelectionParams(
                language_code=language_code,
                name=voice_name
            )
            audio_config = texttospeech.AudioConfig(
                audio_encoding=texttospeech.AudioEncoding.MP3
            )
            synthesis_input = texttospeech.SynthesisInput(text=chunk)
            response = client.synthesize_speech(
                input=synthesis_input,
                voice=voice,
                audio_config=audio_config
            )
            with open(part_file, "wb") as out:
                out.write(response.audio_content)
            print(f"Chunk {i+1} saved to {part_file}")
        except Exception as e:
            print(f"Error synthesizing chunk {i+1}: {e}")
            raise

        # Save state after each chunk
        with state_lock: # Acquire the lock
            completed_chunks.append(part_file)
            save_state({"completed_chunks": completed_chunks, "chunks": chunks})
        progress = (i + 1) / total * 100
        print(f"Progress: {progress:.2f}%")

    # Concurrent synthesis
    with ThreadPoolExecutor(max_workers=num_threads) as executor:
        executor.map(lambda x: synthesize_chunk(*x), remaining_chunks)

    # Merge audio chunks using ffmpeg directly
    import subprocess
    temp_list = f"{base_audio}_filelist.txt"
    with open(temp_list, "w", encoding="utf-8") as f:
        for file in completed_chunks:
            f.write(f"file '{os.path.abspath(file)}'\n")
    if len(completed_chunks) != total:
        print("Not all synthesis tasks completed. Aborting merge.")
        return
    cmd = [
        "ffmpeg",
        "-y",
        "-f", "concat",
        "-safe", "0",
        "-i", temp_list,
        "-c", "copy",
        output_file
    ]
    subprocess.run(cmd, check=True)
    print(f"Merged audio saved to {output_file}")
    os.remove(temp_list)
    # Clean up state file
    try:
        os.remove(state_file)
    except OSError as e:
        print(f"Error removing state file {state_file}: {e}")
    for file in completed_chunks:
        try:
            os.remove(file)
        except OSError as e:
            print(f"Error removing file {file}: {e}")
        text_file = file.replace(output_ext, "-text.txt")
        try:
            os.remove(text_file)
        except OSError as e:
            print(f"Error removing text file {text_file}: {e}")

def main():
    parser = argparse.ArgumentParser(description="Convert text files to speech audio.")
    parser.add_argument("file", help="Path to the input text file")
    parser.add_argument("-f", "--force", action="store_true", default=False,
                        help="Force multi-chunk synthesis without confirmation")
    # Add argument for number of threads
    parser.add_argument("-t", "--threads", type=int, default=1,
                        help="Number of threads for concurrent synthesis (default: 1)")
    args = parser.parse_args()

    filepath = args.file
    if filepath.startswith("http://") or filepath.startswith("https://"):
        print("Extracting text from URL...")
        text = extract_url_text(filepath)
        base_name = sanitize_filename(filepath)
    else:
        text = read_text_file(filepath)
        base_name = os.path.splitext(filepath)[0]

    # Clean entire text to keep only regular speakable characters.
    text = clean_speakable_text(text)

    config = load_config()
    default_provider_key = config.get("default_provider")
    provider_conf = config.get("providers", {}).get(default_provider_key, {})
    provider_type = provider_conf.get("type")
    output_format = provider_conf.get("output_format", "wav")

    # Log current provider information including provider key
    print(f"Using provider key: {default_provider_key}, type: {provider_type}, output file type: {output_format}")

    output_ext = f".{output_format}"
    output_file = base_name + output_ext

    # Output full text sent to TTS engine
    text_output_file = base_name + "-text.txt"
    with open(text_output_file, "w", encoding="utf-8") as f:
        f.write(text)
    print(f"Text saved to {text_output_file}")

    if provider_type == "azure":
        synthesize_azure(text, provider_conf, output_file, force=args.force)
    elif provider_type == "google":
        # Pass num_threads to synthesize_google
        synthesize_google(text, provider_conf, output_file, num_threads=args.threads)
    else:
        print("Provider not supported yet.")

if __name__ == "__main__":
    main()
